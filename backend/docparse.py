import io
import tempfile
import os
from pathlib import Path

MAX_CONTENT = 8000


def parse_document(filename: str, data: bytes) -> dict:
    ext = Path(filename).suffix.lower()
    try:
        # Try markitdown first for supported types
        if ext in (".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".md", ".csv", ".html", ".htm"):
            return _parse_markitdown(filename, data, ext)
        else:
            return {"name": filename, "content": data.decode("utf-8", errors="replace"), "type": "document", "images": []}
    except Exception as e:
        # Fallback to legacy parsers
        try:
            if ext == ".pdf":
                return _parse_pdf(filename, data)
            elif ext == ".docx":
                return _parse_docx(filename, data)
            elif ext in (".xlsx", ".xls"):
                return _parse_xlsx(filename, data)
            elif ext == ".pptx":
                return _parse_pptx(filename, data)
            else:
                return {"name": filename, "content": data.decode("utf-8", errors="replace"), "type": "document", "images": []}
        except Exception as e2:
            return {"name": filename, "content": f"[Fehler beim Lesen: {e2}]", "type": "document", "images": []}


def _parse_markitdown(filename: str, data: bytes, ext: str) -> dict:
    """Use MarkItDown library to convert documents to markdown."""
    try:
        from markitdown import MarkItDown
    except ImportError:
        raise ImportError("markitdown not installed — falling back to legacy parser")

    # Write to temp file because MarkItDown needs a file path
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        md = MarkItDown()
        result = md.convert(tmp_path)
        content = result.text_content if hasattr(result, "text_content") else str(result)
        content = _truncate(content, filename)
        return {"name": filename, "content": content, "type": "document", "images": []}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


def _truncate(text: str, name: str) -> str:
    if len(text) > MAX_CONTENT:
        return text[:MAX_CONTENT] + f"\n\n[… Inhalt gekürzt]"
    return text


def _parse_pdf(filename, data):
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            pages.append(f"[Seite {i + 1}]\n{text}")
    content = "\n\n".join(pages) if pages else "[PDF enthält keinen extrahierbaren Text]"
    return {"name": filename, "content": _truncate(content, filename), "type": "document", "images": []}


def _parse_docx(filename, data):
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append(" | ".join(c.text.strip() for c in row.cells))
        parts.append("\n".join(rows))
    return {"name": filename, "content": _truncate("\n\n".join(parts), filename), "type": "document", "images": []}


def _parse_xlsx(filename, data):
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(c is not None for c in row):
                rows.append(" | ".join("" if c is None else str(c) for c in row))
        if rows:
            sheets.append(f"[Tabelle: {name}]\n" + "\n".join(rows))
    content = "\n\n".join(sheets) if sheets else "[Leere Excel-Datei]"
    return {"name": filename, "content": _truncate(content, filename), "type": "document", "images": []}


def _parse_pptx(filename, data):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if texts:
            slides.append(f"[Folie {i + 1}]\n" + "\n".join(texts))
    content = "\n\n".join(slides) if slides else "[Leere Präsentation]"
    return {"name": filename, "content": _truncate(content, filename), "type": "document", "images": []}
